import fitz
import docx
import os
import time
import requests
import json
import re
import logging
import datetime
from pptx import Presentation
from pymongo import MongoClient

import os
from dotenv import load_dotenv
load_dotenv()

GROQ_API_KEY = os.getenv("GROQ_API_KEY")
MONGO_URI = os.getenv("MONGO_URI", "mongodb://localhost:27017")
client = MongoClient(MONGO_URI)
mongo_db = client.PassionInfotech
history_collection = mongo_db.AI_RISK

# Utility and risk analysis functions

def split_text(text, chunk_size=4000):
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

def extract_text_from_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        text = "\n".join([page.get_text("text") for page in doc])
        doc.close()
        return text
    except Exception as e:
        return None

def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        return None

def extract_text_from_txt(txt_path):
    try:
        with open(txt_path, "r", encoding="utf-8") as file:
            return file.read()
    except Exception as e:
        return None

def extract_text_from_pptx(pptx_path):
    try:
        presentation = Presentation(pptx_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return None

def analyze_risks_with_groq(text):
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }
    chunks = split_text(text)
    risk_reports = []
    for idx, chunk in enumerate(chunks):
        print(f"Analyzing chunk {idx+1}/{len(chunks)}...")
        retry_attempts = 5
        delay = 2
        while retry_attempts > 0:
            prompt = f"""
            You are an AI specializing in risk assessment.
            Given the following document section, analyze potential risks and return ONLY properly formatted JSON.
            IMPORTANT FORMATTING INSTRUCTIONS:
            1. Your response must contain ONLY a single valid JSON object
            2. Do not include any explanatory text before or after the JSON
            3. Do not use markdown code blocks or triple backticks (```) 
            4. Make sure all keys and string values use double quotes, not single quotes
            5. Make sure the JSON syntax is valid - test it carefully
            Use exactly this JSON structure:
            {{
                "RiskID": "RISK-{idx+1:03d}",
                "RiskName": "Brief name of the risk",
                "RiskCategory": "Category such as security, compliance, feasibility, etc.",
                "RiskSeverity": "Low/Medium/High/Critical",
                "RiskDescription": "Detailed description of the identified risk",
                "Probability": "Likelihood of occurrence (Low/Medium/High)",
                "Impact": "Potential impact on the project (Low/Medium/High)",
                "SecurityImplications": "Any security risks associated",
                "TechnicalMitigation": "Specific technical controls, tools, or implementation details to address the risk",
                "NonTechnicalMitigation": "Process changes, training, policies, and organizational measures to address the risk",
                "ContingencyPlan": "Backup plan in case the risk occurs"
            }}
            IMPORTANT NOTES:
            - Ensure a balanced distribution of risks across all severity levels (Low, Medium, High, Critical).
            - Avoid overestimating severity unless justified by the context.
            - Provide specific, actionable technical and non-technical mitigation strategies.
            Document Section:
            {chunk[:3500]}
            """
            payload = {
                "model": "llama-3.3-70b-versatile",
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.4,
                "max_tokens": 1000,
                "response_format": {"type": "json_object"}
            }
            try:
                response = requests.post(
                    "https://api.groq.com/openai/v1/chat/completions",
                    headers=headers,
                    json=payload
                )
                response.raise_for_status()
                content = response.json()["choices"][0]["message"]["content"]
                try:
                    json.loads(content)
                    risk_reports.append(content)
                except json.JSONDecodeError as je:
                    risk_reports.append(json.dumps({
                        "RiskID": f"RISK-ERR-{idx+1:03d}",
                        "RiskName": "API Response Parsing Error",
                        "RiskCategory": "Technical",
                        "RiskSeverity": "Low",
                        "RiskDescription": f"The API response for chunk {idx+1} could not be parsed as valid JSON.",
                        "Probability": "Medium",
                        "Impact": "Low",
                        "SecurityImplications": "None",
                        "TechnicalMitigation": "Review the JSON structure and fixing syntax errors in the API integration code",
                        "NonTechnicalMitigation": "Document this parsing issue and establish a review process for analyzing failed responses",
                        "ContingencyPlan": "Contact support if this error persists"
                    }))
                time.sleep(2)
                break
            except requests.exceptions.HTTPError as e:
                if response.status_code == 429:
                    time.sleep(delay)
                    delay *= 2
                    retry_attempts -= 1
                else:
                    break
            except Exception as e:
                retry_attempts -= 1
                if retry_attempts <= 0:
                    break
    return "\n\n".join(risk_reports) if risk_reports else None

def parse_risk_reports(risk_report_text):
    risk_items = []
    reports = risk_report_text.split("\n\n")
    for idx, report in enumerate(reports):
        try:
            try:
                risk_item = json.loads(report)
                risk_items.append(risk_item)
                continue
            except json.JSONDecodeError:
                pass
            json_text = None
            if "```json" in report:
                json_text = report.split("```json")[1].split("```", 1)[0].strip()
            elif "```" in report:
                json_text = report.split("```", 1)[1].split("```", 1)[0].strip()
            else:
                json_match = re.search(r'(\{.*\})', report, re.DOTALL)
                if json_match:
                    json_text = json_match.group(1).strip()
                else:
                    json_text = report.strip()
            if json_text:
                json_text = json_text.replace("'", '"')
                json_text = re.sub(r'([{,])\s*(\w+):', r'\1 "\2":', json_text)
                json_text = json_text.replace('None', 'null')
                risk_item = json.loads(json_text)
                risk_items.append(risk_item)
            else:
                raise ValueError("No JSON content found in report")
        except (json.JSONDecodeError, ValueError) as e:
            risk_items.append({
                "RiskID": f"parsing-error-{idx+1}",
                "RiskName": "Parsing Error",
                "RiskDescription": f"Failed to parse JSON: {str(e)}.",
                "RiskSeverity": "Unknown",
                "RiskCategory": "Uncategorized",
                "Probability": "Unknown",
                "Impact": "Unknown",
                "SecurityImplications": "Error parsing this risk report",
                "TechnicalMitigation": "Review the raw text manually",
                "NonTechnicalMitigation": "Establish process for handling parsing errors",
                "ContingencyPlan": "Contact support if multiple parsing errors occur"
            })
        except Exception as e:
            risk_items.append({
                "RiskID": f"processing-error-{idx+1}",
                "RiskName": "Processing Error",
                "RiskDescription": f"Error: {str(e)}",
                "RiskSeverity": "Unknown",
                "RiskCategory": "Uncategorized",
                "Probability": "Unknown",
                "Impact": "Unknown",
                "SecurityImplications": "Unknown",
                "TechnicalMitigation": "Review manually",
                "NonTechnicalMitigation": "Document issue",
                "ContingencyPlan": "Retry processing"
            })
    return risk_items

def standardize_severity(severity):
    if not severity:
        return "Unknown"
    severity = severity.lower()
    if severity in ["red"]:
        return "Critical"
    elif severity in ["orange"]:
        return "High"
    elif severity in ["yellow"]:
        return "Medium"
    elif severity in ["green"]:
        return "Low"
    if "critical" in severity:
        return "Critical"
    elif "high" in severity:
        return "High"
    elif "medium" in severity or "med" in severity:
        return "Medium"
    elif "low" in severity:
        return "Low"
    else:
        return "Unknown"

def map_severity(severity):
    if not severity or severity == "Unknown":
        return 5
    severity = str(severity).strip().lower()
    if severity in ["critical", "red"]:
        return 10
    elif severity in ["high", "orange"]:
        return 8
    elif severity in ["medium", "med", "yellow"]:
        return 6
    elif severity in ["low", "green"]:
        return 3
    else:
        return 5

def map_occurrence(probability):
    if not probability or probability == "Unknown":
        return 5
    probability = str(probability).strip().lower()
    if probability in ["high", "very high", "certain"]:
        return 10
    elif probability in ["medium", "moderate", "likely"]:
        return 6
    elif probability in ["low", "unlikely", "rare"]:
        return 3
    else:
        return 5

def calculate_detectability(risk):
    category = str(risk.get("RiskCategory", "")).strip().lower()
    description = str(risk.get("RiskDescription", "")).lower()
    technical_mitigation = str(risk.get("TechnicalMitigation", "")).lower()
    if "security" in category:
        base_score = 7
    elif "technical" in category:
        base_score = 6
    elif "operational" in category:
        base_score = 5
    elif "compliance" in category or "legal" in category:
        base_score = 4
    else:
        base_score = 5
    detection_terms = ["monitor", "alert", "logging", "audit", "detect", "scan", "dashboard", "tracking"]
    control_count = sum(1 for term in detection_terms if term in description or term in technical_mitigation)
    detection_adjustment = min(control_count, 4)
    final_score = max(1, min(10, base_score - detection_adjustment))
    return final_score

def extract_current_controls(risk):
    description = risk.get("RiskDescription", "")
    technical = risk.get("TechnicalMitigation", "")
    non_technical = risk.get("NonTechnicalMitigation", "")
    control_indicators = [
        "currently", "existing", "in place", "implemented", 
        "already", "present", "established"
    ]
    controls = []
    for text in [description, technical, non_technical]:
        sentences = text.split(". ")
        for sentence in sentences:
            if any(indicator in sentence.lower() for indicator in control_indicators):
                controls.append(sentence.strip())
    if not controls:
        return "No existing controls documented."
    return " ".join(controls)

def parse_suggested_actions(suggested_fix):
    if not suggested_fix or suggested_fix == "No immediate action required.":
        return []
    actions = []
    lines = suggested_fix.split('\n')
    current_action = ""
    for line in lines:
        line = line.strip()
        if not line:
            if current_action:
                actions.append(current_action)
                current_action = ""
            continue
        if (line.startswith('- ') or line.startswith('• ') or 
            any(line.startswith(f"{i}.") for i in range(1, 11))):
            if current_action:
                actions.append(current_action)
            current_action = line
        else:
            if current_action:
                current_action += " " + line
            else:
                current_action = line
    if current_action:
        actions.append(current_action)
    if not actions and suggested_fix:
        actions = [suggested_fix]
    return actions


# Enhanced FMEA logic with tiered thresholds and action levels
def calculate_rpn_and_suggest_fixes(risk_items):
    RPN_HIGH = 200
    RPN_MODERATE = 100
    CN_HIGH = 70  # Severity (1-10) × Occurrence threshold
    fmea_results = []
    severity_distribution = {"Critical": 0, "High": 0, "Medium": 0, "Low": 0}
    for risk in risk_items:
        try:
            severity = map_severity(risk.get("RiskSeverity", "Low"))
            occurrence = map_occurrence(risk.get("Probability", "Low"))
            detectability = calculate_detectability(risk)
            rpn = severity * occurrence * detectability
            cn = severity * occurrence
            # Action level logic
            if rpn >= RPN_HIGH or cn >= CN_HIGH or severity >= 9:
                action_level = "Immediate"
            elif rpn >= RPN_MODERATE or cn >= (CN_HIGH // 2):
                action_level = "Preventive"
            elif severity >= 8:
                action_level = "ManualReview"
            else:
                action_level = "Monitor"
            # Logging for debugging
            print(f"Risk: {risk.get('RiskName')} - RPN: {rpn} (S:{severity}, O:{occurrence}, D:{detectability}), CN: {cn}, ActionLevel: {action_level}")
            risk["FMEA"] = {
                "Severity": severity,
                "Occurrence": occurrence,
                "Detection": detectability,
                "RPN": rpn,
                "CurrentControls": extract_current_controls(risk),
                "RecommendedActions": [],
                "ActionStatus": "Not Started",
                "ResponsiblePerson": "",
                "TargetDate": "",
                "ActionTaken": "",
                "UpdatedRPN": None,
                "ActionLevel": action_level
            }
            risk["RPN"] = rpn
            risk["ActionLevel"] = action_level
            # Generate suggestions only for Immediate/Preventive
            if action_level in ["Immediate", "Preventive"]:
                suggested_actions = generate_ai_suggestions(risk)
                risk["FMEA"]["RecommendedActions"] = parse_suggested_actions(suggested_actions)
                risk["SuggestedFix"] = suggested_actions
            elif action_level == "ManualReview":
                risk["SuggestedFix"] = "Manual review required by risk team."
            else:
                risk["SuggestedFix"] = "Monitor as part of regular review."
            severity_level = risk.get("RiskSeverity", "Unknown")
            if severity_level in severity_distribution:
                severity_distribution[severity_level] += 1
            fmea_results.append(risk)
        except Exception as e:
            risk["RPN"] = 0
            risk["SuggestedFix"] = f"Error calculating RPN: {e}"
            fmea_results.append(risk)
    return fmea_results

def generate_ai_suggestions(risk):
    try:
        prompt = f"""
        You are an AI specializing in risk mitigation strategies.
        Given the following risk details, provide specific technical and non-technical mitigation strategies:
        Risk Details:
        - Risk Name: {risk.get('RiskName', 'Unnamed Risk')}
        - Risk Category: {risk.get('RiskCategory', 'Uncategorized')}
        - Risk Severity: {risk.get('RiskSeverity', 'Unknown')}
        - Probability: {risk.get('Probability', 'Unknown')}
        - Impact: {risk.get('Impact', 'Unknown')}
        - Risk Description: {risk.get('RiskDescription', 'No description provided')}
        Provide actionable and detailed mitigation strategies.
        """
        headers = {
            "Authorization": f"Bearer {GROQ_API_KEY}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": "llama-3.3-70b-versatile",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.4,
            "max_tokens": 500
        }
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers=headers,
            json=payload
        )
        response.raise_for_status()
        content = response.json()["choices"][0]["message"]["content"]
        return content.strip()
    except Exception as e:
        return "Error generating AI suggestions. Please review the risk manually."

def calculate_overall_risk(risk_items):
    risk_levels = {"Critical": 4, "High": 3, "Medium": 2, "Low": 1}
    max_risk_level = 0
    summary = []
    for risk in risk_items:
        severity = risk.get("RiskSeverity", "Low")
        max_risk_level = max(max_risk_level, risk_levels.get(severity, 1))
        summary.append(f"{risk.get('RiskName', 'Unnamed Risk')}: {severity}")
    overall_level = next((key for key, value in risk_levels.items() if value == max_risk_level), "Low")
    return overall_level, summary
